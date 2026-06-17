from __future__ import annotations

import csv
import json
import math
import re
import subprocess
from collections import Counter
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PRESENTATION_DIR = ROOT / "presentation"
DATA_DIR = PRESENTATION_DIR / "data"
OUTPUT_PATH = PRESENTATION_DIR / "pg_case_factory_project_report_20260415.pptx"
VALIDATION_CACHE = DATA_DIR / "runtime_execution_validation_20260415.json"

WIDE = 13.333
HIGH = 7.5

NAVY = RGBColor(17, 31, 59)
BLUE = RGBColor(37, 99, 235)
SKY = RGBColor(125, 211, 252)
TEAL = RGBColor(13, 148, 136)
AMBER = RGBColor(245, 158, 11)
RED = RGBColor(220, 38, 38)
GREEN = RGBColor(34, 197, 94)
INK = RGBColor(15, 23, 42)
MUTED = RGBColor(71, 85, 105)
LIGHT = RGBColor(248, 250, 252)
PANEL = RGBColor(226, 232, 240)
WHITE = RGBColor(255, 255, 255)

TITLE_FONT = "Aptos Display"
BODY_FONT = "PingFang SC"
MONO_FONT = "Menlo"


@dataclass(frozen=True)
class ExecutionValidation:
    total_files: int
    script_success_count: int
    script_failure_count: int
    verify_hit_1_count: int
    verify_zero_row_count: int
    error_counts: dict[str, int]
    error_factor_pairs: dict[str, dict[str, int]]
    sample_failures: list[dict[str, str]]

    def to_record(self) -> dict[str, object]:
        return {
            "total_files": self.total_files,
            "script_success_count": self.script_success_count,
            "script_failure_count": self.script_failure_count,
            "verify_hit_1_count": self.verify_hit_1_count,
            "verify_zero_row_count": self.verify_zero_row_count,
            "error_counts": self.error_counts,
            "error_factor_pairs": self.error_factor_pairs,
            "sample_failures": self.sample_failures,
        }


def read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8")


def read_json(path: Path) -> dict:
    return json.loads(read_text(path))


def read_yaml(path: Path) -> dict:
    return yaml.safe_load(read_text(path))


def yaml_block_from_skill(raw_text: str) -> dict:
    match = re.search(r"```yaml\s*(.*?)```", raw_text, re.DOTALL)
    if not match:
        raise ValueError("skill file missing YAML block")
    return yaml.safe_load(match.group(1))


def clip(text: str, limit: int) -> str:
    compact = " ".join(text.strip().split())
    return compact if len(compact) <= limit else compact[: limit - 1] + "…"


def make_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(WIDE)
    prs.slide_height = Inches(HIGH)
    return prs


def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, text: str, dark: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(0.35), Inches(7.1), Inches(12.6), Inches(0.22))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = BODY_FONT
    run.font.size = Pt(9)
    run.font.color.rgb = PANEL if dark else MUTED
    p.alignment = PP_ALIGN.RIGHT


def add_title(slide, title: str, subtitle: str | None = None, dark: bool = False) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.38), Inches(11.9), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = TITLE_FONT
    run.font.bold = True
    run.font.size = Pt(26)
    run.font.color.rgb = WHITE if dark else NAVY
    if subtitle:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.name = BODY_FONT
        run2.font.size = Pt(12)
        run2.font.color.rgb = PANEL if dark else MUTED
        p2.space_before = Pt(5)


def add_box(slide, left: float, top: float, width: float, height: float, fill_color: RGBColor, line_color: RGBColor | None = None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color or fill_color
    return shape


def fill_shape_text(shape, text: str, font_size: int = 16, color: RGBColor = INK, bold: bool = False, font_name: str = BODY_FONT, valign=MSO_ANCHOR.TOP, margin: float = 0.1) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_bullet_block(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    bullets: Iterable[str],
    title: str | None = None,
    fill_color: RGBColor = WHITE,
    title_color: RGBColor = NAVY,
    body_color: RGBColor = INK,
    line_color: RGBColor = PANEL,
) -> None:
    shape = add_box(slide, left, top, width, height, fill_color, line_color)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
    if title:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.name = BODY_FONT
        run.font.bold = True
        run.font.size = Pt(16)
        run.font.color.rgb = title_color
        p.space_after = Pt(4)
    else:
        tf.paragraphs[0].text = ""
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.level = 0
        p.font.name = BODY_FONT
        p.font.size = Pt(12)
        p.font.color.rgb = body_color
        p.space_after = Pt(2)


def add_metric_card(slide, left: float, top: float, width: float, height: float, value: str, label: str, accent: RGBColor) -> None:
    shape = add_box(slide, left, top, width, height, WHITE, PANEL)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = value
    r1.font.name = TITLE_FONT
    r1.font.bold = True
    r1.font.size = Pt(24)
    r1.font.color.rgb = accent
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = label
    r2.font.name = BODY_FONT
    r2.font.size = Pt(11)
    r2.font.color.rgb = MUTED
    strip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(0.12),
        Inches(height),
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.color.rgb = accent


def add_code_block(slide, left: float, top: float, width: float, height: float, text: str, title: str | None = None) -> None:
    shape = add_box(slide, left, top, width, height, RGBColor(241, 245, 249), RGBColor(203, 213, 225))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    if title:
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = BODY_FONT
        r.font.bold = True
        r.font.size = Pt(13)
        r.font.color.rgb = NAVY
    for line in text.splitlines():
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.name = MONO_FONT
        r.font.size = Pt(10)
        r.font.color.rgb = INK
        p.space_after = Pt(0)


def add_table_slide(slide, left: float, top: float, width: float, height: float, columns: list[str], rows: list[list[str]], header_fill: RGBColor = NAVY) -> None:
    table = slide.shapes.add_table(len(rows) + 1, len(columns), Inches(left), Inches(top), Inches(width), Inches(height)).table
    for idx, name in enumerate(columns):
        cell = table.cell(0, idx)
        cell.text = name
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.name = BODY_FONT
                run.font.bold = True
                run.font.size = Pt(10)
                run.font.color.rgb = WHITE
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 else RGBColor(248, 250, 252)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = BODY_FONT
                    run.font.size = Pt(10)
                    run.font.color.rgb = INK


def execution_validation() -> ExecutionValidation:
    if VALIDATION_CACHE.exists():
        cached = read_json(VALIDATION_CACHE)
        return ExecutionValidation(
            total_files=int(cached["total_files"]),
            script_success_count=int(cached["script_success_count"]),
            script_failure_count=int(cached["script_failure_count"]),
            verify_hit_1_count=int(cached["verify_hit_1_count"]),
            verify_zero_row_count=int(cached["verify_zero_row_count"]),
            error_counts={str(k): int(v) for k, v in dict(cached["error_counts"]).items()},
            error_factor_pairs={str(k): {str(k2): int(v2) for k2, v2 in dict(v).items()} for k, v in dict(cached["error_factor_pairs"]).items()},
            sample_failures=list(cached["sample_failures"]),
        )

    sql_dir = ROOT / "artifacts" / "generated_sql" / "normal_table-create_index"
    binding_path = ROOT / "artifacts" / "intermediates" / "normal_table-create_index-bindings.tsv"
    with binding_path.open(encoding="utf-8", newline="") as handle:
        rows = list(csv.DictReader(handle, delimiter="\t"))
    by_file = {row["sql_file"]: row for row in rows}

    total_files = 0
    script_success_count = 0
    verify_hit_1_count = 0
    verify_zero_row_count = 0
    error_counts: Counter[str] = Counter()
    error_factor_pairs: dict[str, Counter[str]] = {}
    sample_failures: list[dict[str, str]] = []

    for path in sorted(sql_dir.glob("*.sql")):
        total_files += 1
        proc = subprocess.run(
            ["psql", "-d", "postgres", "-v", "ON_ERROR_STOP=1", "-f", str(path)],
            capture_output=True,
            text=True,
            check=False,
        )
        merged = (proc.stdout or "") + (proc.stderr or "")
        if proc.returncode == 0:
            script_success_count += 1
        if "(1 row)" in merged:
            verify_hit_1_count += 1
        if "(0 rows)" in merged:
            verify_zero_row_count += 1
        if proc.returncode == 0:
            continue

        error_message = ""
        for line in reversed(merged.splitlines()):
            if "ERROR:" in line:
                error_message = line.split("ERROR:", 1)[1].strip()
                break
        error_counts[error_message] += 1

        row = by_file.get(str(path.resolve())) or by_file.get(str(path))
        if row:
            pair_key = f"{row['indexed_type_family']} | {row['predicate_mode']}"
            error_factor_pairs.setdefault(error_message, Counter())[pair_key] += 1

        if len(sample_failures) < 6:
            sample_failures.append(
                {
                    "file": path.name,
                    "error": error_message,
                    "tail": "\n".join(merged.splitlines()[-5:]),
                }
            )

    record = ExecutionValidation(
        total_files=total_files,
        script_success_count=script_success_count,
        script_failure_count=total_files - script_success_count,
        verify_hit_1_count=verify_hit_1_count,
        verify_zero_row_count=verify_zero_row_count,
        error_counts=dict(error_counts),
        error_factor_pairs={k: dict(v) for k, v in error_factor_pairs.items()},
        sample_failures=sample_failures,
    )
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    VALIDATION_CACHE.write_text(json.dumps(record.to_record(), ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return record


def build_deck() -> None:
    PRESENTATION_DIR.mkdir(parents=True, exist_ok=True)
    DATA_DIR.mkdir(parents=True, exist_ok=True)

    readme = read_text(ROOT / "README.md")
    project_structure = read_text(ROOT / "PROJECT_STRUCTURE.md")
    evaluation = read_json(ROOT / "artifacts" / "evaluations" / "TASK_20260414_120051.json")
    normalized = read_yaml(ROOT / "artifacts" / "intermediates" / "TASK_20260414_120051-normalized.yaml")
    abstract_case = read_yaml(ROOT / "artifacts" / "intermediates" / "TASK_20260414_120051-abstract.yaml")
    manifest = read_json(ROOT / "artifacts" / "intermediates" / "TASK_20260414_120051-manifest.json")
    generator_program = read_text(ROOT / "artifacts" / "generated_programs" / "normal_table-create_index.py")
    sample_sql = read_text(ROOT / "artifacts" / "generated_sql" / "normal_table-create_index" / "create_index_AC_CREATE_INDEX_000001_0002.sql")
    object_template = read_text(ROOT / "objects" / "table_02.sql" / "normal_table" / "all_column_types.sql")
    create_index_skill = read_text(ROOT / "skills" / "index" / "create_index.skill")

    using_superpowers = read_text(Path("/Users/yuyu/.codex/superpowers/skills/using-superpowers/SKILL.md"))
    brainstorming = read_text(Path("/Users/yuyu/.codex/superpowers/skills/brainstorming/SKILL.md"))
    writing_plans = read_text(Path("/Users/yuyu/.codex/superpowers/skills/writing-plans/SKILL.md"))
    orchestrator_skill = read_text(Path("/Users/yuyu/.codex/skills/pg-case-factory-orchestrator/SKILL.md"))
    auditor_skill = read_text(Path("/Users/yuyu/.codex/skills/pg-case-factory-plan-auditor/SKILL.md"))
    superpowers_readme = read_text(Path("/Users/yuyu/.codex/superpowers/README.md"))

    validation = execution_validation()

    structured_config = yaml_block_from_skill(create_index_skill)["structured_config"]
    factors = structured_config["factors"]
    important_factors = structured_config["coverage_policy"]["main_combination_axes"]
    non_important_factors = structured_config["coverage_policy"]["non_main_factors"]

    factor_rows = []
    for name, doc in factors.items():
        values = doc.get("values") or []
        factor_rows.append(
            [
                name,
                str(doc.get("label", name)),
                str(doc.get("importance", "non_important")),
                str(len(values)),
                ", ".join(str(item["key"] if isinstance(item, dict) else item) for item in values[:6])
                + (" ..." if len(values) > 6 else ""),
            ]
        )

    with (ROOT / "artifacts" / "intermediates" / "normal_table-create_index-bindings.tsv").open(encoding="utf-8", newline="") as handle:
        binding_rows = list(csv.DictReader(handle, delimiter="\t"))
    factor_distribution: dict[str, Counter[str]] = {}
    for field in (
        "indexed_type_family",
        "uniqueness",
        "order",
        "nulls_position",
        "predicate_mode",
        "index_name_mode",
        "concurrently",
        "if_not_exists",
        "include_columns",
        "storage_parameter_mode",
        "tablespace_mode",
    ):
        factor_distribution[field] = Counter(row[field] for row in binding_rows)

    prs = make_presentation()

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_title(
        slide,
        "pg_case_factory 项目复盘与提示词体系分析",
        "基于仓库实物、Codex 本地 skill 文件、现有 artifacts，以及 2026-04-15 的 PostgreSQL 16.13 实跑验证生成",
        dark=True,
    )
    hero = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.0), Inches(2.0))
    tf = hero.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "一个围绕 PostgreSQL CREATE INDEX 的 SQL 用例工厂 demo。"
    r.font.name = TITLE_FONT
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "核心结论：生成 1188 个 SQL，用仓库内审计视角看是“覆盖通过”；拿真实 PostgreSQL 执行，864 个脚本成功、324 个失败。"
    r2.font.name = BODY_FONT
    r2.font.size = Pt(18)
    r2.font.color.rgb = RGBColor(219, 234, 254)
    add_metric_card(slide, 0.8, 4.25, 2.3, 1.25, "1188", "生成 SQL 数", SKY)
    add_metric_card(slide, 3.35, 4.25, 2.3, 1.25, "1", "抽象用例数", GREEN)
    add_metric_card(slide, 5.9, 4.25, 2.3, 1.25, "864", "实跑成功脚本", TEAL)
    add_metric_card(slide, 8.45, 4.25, 2.3, 1.25, "324", "实跑失败脚本", AMBER)
    add_metric_card(slide, 11.0, 4.25, 1.5, 1.25, "863", "验证命中", BLUE)
    add_footer(slide, f"Workspace: {ROOT}", dark=True)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "1. 项目定位与结论摘要", "结论页先给管理层和技术评审一个统一视图")
    add_bullet_block(
        slide,
        0.55,
        1.25,
        6.0,
        4.95,
        [
            "仓库定位是 SQL-first PostgreSQL 16.4 DDL test case generator demo，不是完整的测试执行平台。",
            "当前自然语言入口只支持一句话意图：'覆盖所有列类型和创建索引语句'，会固化到 normal_table + create_index。",
            "主流程由 Python 代码驱动，领域知识不在模型权重里，而在 skills/*.skill 的 YAML 配置里。",
            "repo 内建的 audit 只证明“覆盖空间是否被枚举出来”，不证明 SQL 在真实数据库里一定能执行成功。",
            "本次补做的 PostgreSQL 16.13 实跑验证发现：覆盖通过 ≠ 实库通过，失败集中暴露出 skill 渲染规则的语义问题。",
        ],
        title="一句话总结",
    )
    add_bullet_block(
        slide,
        6.8,
        1.25,
        5.95,
        4.95,
        [
            "仓库内审计：1 个抽象 aggregate case、1188 个 SQL、8 个非重要因子全部覆盖、audit_result.passed=true。",
            "实库验证：1188 总数，864 成功，324 失败，863 个验证查询命中 1 行，1 个自动命名场景返回 0 行。",
            "两类主错误：216 个 'functions in index predicate must be marked IMMUTABLE'；108 个 'point has no default operator class for btree'。",
            "计划层与执行层还有一个落差：test plan 写了两条 lifecycle row，但实际 SQL 脚本只实现了 drop_table 路径。",
        ],
        title="关键数字",
    )
    add_footer(slide, "Sources: README.md, PROJECT_STRUCTURE.md, artifacts/evaluations/TASK_20260414_120051.json")

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "2. 仓库怎么组织", "项目代码、规则、模板和产物被拆成四个主区")
    tree_text = """pg_case_factory/
├─ README.md / PROJECT_STRUCTURE.md
├─ objects/
│  └─ table_02.sql/normal_table/all_column_types.sql
├─ skills/
│  ├─ common/{output_script_style, common_factor_policy, naming_rules}.skill
│  └─ index/create_index.skill
├─ src/pg_case_factory/
│  ├─ request_parser.py / task_normalizer.py / skill_loader.py
│  ├─ abstract_case_generator.py / sql_renderer.py / evaluator.py
│  └─ workflow.py / cli.py / models.py / artifact_store.py
└─ artifacts/
   ├─ generated_programs/ 1
   ├─ generated_sql/ 1188
   ├─ test_plans/ 1
   ├─ evaluations/ 2
   └─ intermediates/ 6"""
    add_code_block(slide, 0.55, 1.25, 6.2, 5.5, tree_text, title="仓库骨架")
    add_bullet_block(
        slide,
        6.95,
        1.25,
        5.8,
        5.5,
        [
            "`objects/` 只负责基础对象模板，当前只有一张“全列类型普通表”。",
            "`skills/` 负责提示词式规则配置：既描述覆盖空间，也描述 SQL 片段映射和渲染模板。",
            "`src/pg_case_factory/` 负责把请求、skill、对象模板、审计逻辑串成主流程。",
            "`artifacts/` 是唯一运行产物根目录，每次运行前会被 workflow 清空重建。",
            "我把最终 PPT 放到 presentation/，避免下次 handle_chat_request 运行时被 artifacts 清掉。",
        ],
        title="设计取向",
    )
    add_footer(slide, "Sources: PROJECT_STRUCTURE.md, src/pg_case_factory/artifact_store.py")

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "3. 端到端主流程", "代码实现对应的是一条固定 pipeline，而不是 agent 自由发挥")
    flow_boxes = [
        (0.6, "用户请求", "parse_request\n只识别 CREATE INDEX + 所有列类型"),
        (2.2, "skill 装载", "load_skill\n读取 create_index.skill 中的 YAML DSL"),
        (3.8, "任务归一化", "normalize_task\n绑定对象模板、因子、默认值、规则路径"),
        (5.4, "抽象展开", "generate_abstract_cases + build_atomic_bindings"),
        (7.0, "SQL 渲染", "render_sql_script\n对象模板 + 片段绑定 + 命名规则"),
        (8.6, "产物落盘", "plan/program/sql/manifest/evaluation"),
        (10.2, "覆盖审计", "build_coverage_report\n只审 coverage，不连数据库"),
    ]
    for left, title, body in flow_boxes:
        box = add_box(slide, left, 2.0, 1.35, 2.0, WHITE, PANEL)
        fill_shape_text(box, f"{title}\n{body}", font_size=12, color=INK)
    for idx in range(len(flow_boxes) - 1):
        arrow = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.CHEVRON,
            Inches(1.85 + idx * 1.6),
            Inches(2.65),
            Inches(0.25),
            Inches(0.45),
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = SKY
        arrow.line.color.rgb = SKY
    add_bullet_block(
        slide,
        0.8,
        4.55,
        12.0,
        1.7,
        [
            "workflow.py 是真正的编排中心，负责创建 run_id、准备 artifacts、生成 SQL 文件、写 manifest 和评估 JSON/Markdown。",
            "CLI 只是薄壳：`PYTHONPATH=src python3 -m pg_case_factory demo --request \"覆盖所有列类型和创建索引语句\"`。",
        ],
        title="实现要点",
    )
    add_footer(slide, "Sources: src/pg_case_factory/workflow.py, cli.py, request_parser.py")

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "4. 请求如何被解析和约束", "当前 demo 的入口非常窄，但可控")
    add_code_block(
        slide,
        0.55,
        1.25,
        5.8,
        2.45,
        "\n".join(read_text(ROOT / "src" / "pg_case_factory" / "request_parser.py").splitlines()[0:26]),
        title="request_parser.py 关键逻辑",
    )
    add_bullet_block(
        slide,
        6.6,
        1.25,
        6.1,
        2.45,
        [
            "请求文本不能为空。",
            "必须显式包含 'create index' 或 '创建索引'。",
            "必须显式包含 '所有列类型' 或 'all column types'。",
            "否则直接抛 ValueError；当前 demo 不支持其他语句、对象或覆盖范围。",
        ],
        title="输入门槛",
    )
    add_bullet_block(
        slide,
        0.55,
        3.95,
        12.15,
        2.3,
        [
            f"标准化结果 task_slug = {normalized['task_slug']}，base_object = {normalized['base_object_key']}，statement = {normalized['target_statement']}。",
            f"important_factors = {', '.join(normalized['important_factors'])}。",
            f"non_important_factors = {', '.join(normalized['non_important_factors'])}。",
            "python_expand_threshold = 200，而 estimated_sql_count = 1188，所以会保留生成程序文件。",
        ],
        title="归一化后固定下来的运行配置",
    )
    add_footer(slide, "Sources: src/pg_case_factory/request_parser.py, task_normalizer.py, artifacts/intermediates/*-normalized.yaml")

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_title(slide, "5. 可见 Codex 提示词栈怎么设计", "这里只讲本机可读取到的规则层，不直接复述隐藏运行时 prompt 原文", dark=True)
    add_bullet_block(
        slide,
        0.55,
        1.2,
        3.0,
        4.9,
        [
            "运行时基础指令层：角色、沟通风格、编辑约束、验证要求、协作模式。",
            "我在本机 session 归档里能看到这层的结构痕迹，但 PPT 只做摘要，不整段曝光内部提示原文。",
        ],
        title="L0 运行时基座",
        fill_color=RGBColor(30, 41, 59),
        title_color=WHITE,
        body_color=RGBColor(226, 232, 240),
        line_color=RGBColor(51, 65, 85),
    )
    add_bullet_block(
        slide,
        3.75,
        1.2,
        2.9,
        4.9,
        [
            "using-superpowers：要求每次任务先检查 skill。",
            "brainstorming：先理解上下文、问清楚、给设计。",
            "writing-plans：把设计拆成可执行计划。",
        ],
        title="L1 通用流程 skill",
        fill_color=RGBColor(30, 41, 59),
        title_color=WHITE,
        body_color=RGBColor(226, 232, 240),
        line_color=RGBColor(51, 65, 85),
    )
    add_bullet_block(
        slide,
        6.85,
        1.2,
        2.9,
        4.9,
        [
            "pg-case-factory-orchestrator：从 repo discovery 生成 lifecycle plan TSV。",
            "pg-case-factory-plan-auditor：重算期望 row，审 plan 是否缺失或冗余。",
        ],
        title="L2 仓库专属 skill",
        fill_color=RGBColor(30, 41, 59),
        title_color=WHITE,
        body_color=RGBColor(226, 232, 240),
        line_color=RGBColor(51, 65, 85),
    )
    add_bullet_block(
        slide,
        9.95,
        1.2,
        2.8,
        4.9,
        [
            "create_index.skill：领域规则、因子、模板、映射。",
            "common_factor_policy.skill：覆盖策略。",
            "output_script_style.skill：输出脚本风格。",
            "naming_rules.skill：命名约束。",
        ],
        title="L3 领域 DSL",
        fill_color=RGBColor(30, 41, 59),
        title_color=WHITE,
        body_color=RGBColor(226, 232, 240),
        line_color=RGBColor(51, 65, 85),
    )
    add_footer(slide, "Sources: ~/.codex/superpowers/skills/*.md, ~/.codex/skills/pg-case-factory-*.md, skills/*.skill", dark=True)

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "6. 关键提示词文件分别约束什么", "提示词设计不是一段大 prompt，而是多份职责清晰的规则文件")
    add_table_slide(
        slide,
        0.45,
        1.15,
        12.45,
        5.7,
        ["文件", "主要职责", "对本项目的影响"],
        [
            ["using-superpowers", "先检查 skill，再行动", "保证任务先走方法论，不直接开写"],
            ["brainstorming", "先做上下文探索和设计", "适合讲方案、做 PPT、大改功能前定结构"],
            ["writing-plans", "把设计拆成细粒度计划", "把 spec 变成可执行步骤"],
            ["pg-case-factory-orchestrator", "从 repo discovery 生成 plan TSV", "把自然语言请求映射到对象 + statement + lifecycle rows"],
            ["pg-case-factory-plan-auditor", "只审 plan 完整性", "保证 plan 行与 repo discovery 对齐"],
            ["create_index.skill", "因子模型 + SQL 模板", "真正决定会生成哪些 case、生成什么 SQL"],
        ],
    )
    add_footer(
        slide,
        "Sources: ~/.codex/superpowers/skills/using-superpowers/SKILL.md, brainstorming/SKILL.md, writing-plans/SKILL.md, ~/.codex/skills/pg-case-factory-*.md, skills/index/create_index.skill",
    )

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "7. create_index.skill 是项目真正的“领域提示词”", "这里既定义覆盖空间，也定义 SQL 怎么拼出来")
    add_bullet_block(
        slide,
        0.55,
        1.15,
        5.8,
        5.55,
        [
            f"skill_name = {structured_config['skill_name']}，statement = {structured_config['statement']['name']}。",
            "T1：对象与方法；T2：主语义；T3：挂靠因子。",
            f"重要因子 {len(important_factors)} 个：{', '.join(important_factors)}。",
            f"非重要因子 {len(non_important_factors)} 个：{', '.join(non_important_factors)}。",
            "rendering.factor_value_bindings 把每个因子值映射成 SQL 片段，占位符再套进 statement_template。",
        ],
        title="结构",
    )
    add_table_slide(
        slide,
        6.55,
        1.15,
        6.25,
        5.55,
        ["因子", "标签", "重要性", "取值数", "前几个取值"],
        factor_rows[:8],
    )
    add_footer(slide, "Sources: skills/index/create_index.skill")

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "8. 用例是怎么被展开出来的", "抽象 case 只有 1 个，但会膨胀成 1188 条原子 binding")
    add_metric_card(slide, 0.7, 1.25, 2.2, 1.1, "11", "indexed_type_family", SKY)
    add_metric_card(slide, 3.05, 1.25, 1.8, 1.1, "2", "uniqueness", SKY)
    add_metric_card(slide, 5.0, 1.25, 1.8, 1.1, "3", "order", SKY)
    add_metric_card(slide, 6.95, 1.25, 1.8, 1.1, "3", "nulls_position", SKY)
    add_metric_card(slide, 8.9, 1.25, 1.8, 1.1, "6", "predicate_mode", SKY)
    add_metric_card(slide, 10.85, 1.25, 1.8, 1.1, "1188", "笛卡尔积总数", GREEN)
    add_bullet_block(
        slide,
        0.55,
        2.7,
        6.0,
        3.7,
        [
            "abstract_case_generator.build_atomic_bindings 先对 5 个重要因子做完整笛卡尔积。",
            "公式：11 × 2 × 3 × 3 × 6 = 1188。",
            "再把 8 个非重要因子的非默认值按 rotate_attach 策略轮转挂靠到前面的主组合上。",
            "本次 main_count 足够大，所以 clone_skeleton_used = false，没有复制主骨架。",
            "estimated_sql_count = 1188 > python_expand_threshold 200，因此 python_expansion_used = true。",
        ],
        title="展开算法",
    )
    add_bullet_block(
        slide,
        6.8,
        2.7,
        6.0,
        3.7,
        [
            "这套 rotate_attach 在当前数据上有一个实际效果：非重要因子的绝大多数非默认值只出现 1 次。",
            "例子：index_name_mode 里 implicit_generated_name 只出现 1 次；concurrently=on 只出现 1 次。",
            "所以它更像“确保值至少覆盖一次”，不是“均匀混合”或“组合增强”。",
        ],
        title="策略副作用",
    )
    add_footer(slide, "Sources: src/pg_case_factory/abstract_case_generator.py, artifacts/intermediates/*-abstract.yaml, bindings.tsv")

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "9. 模板与渲染", "对象模板负责列类型全量覆盖，渲染层负责把 binding 变成 SQL")
    add_code_block(
        slide,
        0.5,
        1.2,
        6.15,
        5.55,
        "\n".join(object_template.splitlines()[:22]),
        title="objects/table_02.sql/normal_table/all_column_types.sql（节选）",
    )
    add_code_block(
        slide,
        6.85,
        1.2,
        5.95,
        5.55,
        "\n".join(sample_sql.splitlines()[70:89]),
        title="生成 SQL 样例（create_index_0002，节选）",
    )
    add_footer(slide, "Sources: objects/table_02.sql/normal_table/all_column_types.sql, artifacts/generated_sql/...0002.sql, sql_renderer.py")

    # Slide 11
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "10. 运行后到底产出了什么", "artifacts 里每类文件的职责都很清晰")
    add_metric_card(slide, 0.6, 1.2, 2.1, 1.0, "1", "generated_programs", BLUE)
    add_metric_card(slide, 2.95, 1.2, 2.1, 1.0, "1188", "generated_sql", GREEN)
    add_metric_card(slide, 5.3, 1.2, 2.1, 1.0, "1", "test_plans", SKY)
    add_metric_card(slide, 7.65, 1.2, 2.1, 1.0, "2", "evaluations", TEAL)
    add_metric_card(slide, 10.0, 1.2, 2.1, 1.0, "6", "intermediates", AMBER)
    add_bullet_block(
        slide,
        0.55,
        2.55,
        6.05,
        3.95,
        [
            "test plan TSV 有 2 条 lifecycle rows：",
            "1) create_table>create_index>verify_index>drop_table",
            "2) create_table>create_index>verify_index>drop_index>drop_table",
            "但 binding TSV 和最终 SQL 文件只落成了前一种 operation_chain。",
            "manifest.json 记录了全部 1188 个 SQL 的绝对路径和所有中间产物路径。",
        ],
        title="计划层与执行层",
    )
    add_code_block(
        slide,
        6.85,
        2.55,
        5.95,
        3.95,
        "\n".join(generator_program.splitlines()[:18]),
        title="generated_programs/normal_table-create_index.py",
    )
    add_footer(slide, "Sources: artifacts/test_plans/normal_table-create_index.tsv, artifacts/intermediates/*manifest.json, generated_programs/*.py")

    # Slide 12
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "11. 仓库内置审计怎么看这次生成", "这是 repo 原生评估结果，不是我额外加的实库执行")
    add_bullet_block(
        slide,
        0.55,
        1.2,
        6.0,
        5.45,
        [
            f"total_abstract_cases = {evaluation['summary']['total_abstract_cases']}",
            f"total_generated_sql_files = {evaluation['summary']['total_generated_sql_files']}",
            f"clone_skeleton_used = {evaluation['summary']['clone_skeleton_used']}",
            f"python_expansion_used = {evaluation['summary']['python_expansion_used']}",
            f"audit_result.passed = {evaluation['audit_result']['passed']}",
            "重要因子交叉数量 theoretical = 1188，actual_generated_main_combinations = 1188。",
            "8 个非重要因子全部显示 fully_covered = true。",
        ],
        title="JSON 评估摘要",
    )
    md_summary = "\n".join(read_text(ROOT / "artifacts" / "evaluations" / "TASK_20260414_120051.md").splitlines()[:18])
    add_code_block(slide, 6.8, 1.2, 6.0, 5.45, md_summary, title="中文评估 Markdown（节选）")
    add_footer(slide, "Sources: artifacts/evaluations/TASK_20260414_120051.json/.md")

    # Slide 13
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_title(slide, "12. 我补做的真实 PostgreSQL 执行验证", "环境：psql 16.13，pg_isready = /tmp:5432 accepting connections，目标库 = postgres", dark=True)
    add_metric_card(slide, 0.8, 1.45, 2.2, 1.15, str(validation.total_files), "总 SQL 文件", SKY)
    add_metric_card(slide, 3.2, 1.45, 2.2, 1.15, str(validation.script_success_count), "脚本执行成功", GREEN)
    add_metric_card(slide, 5.6, 1.45, 2.2, 1.15, str(validation.script_failure_count), "脚本执行失败", AMBER)
    add_metric_card(slide, 8.0, 1.45, 2.2, 1.15, str(validation.verify_hit_1_count), "验证命中 1 行", TEAL)
    add_metric_card(slide, 10.4, 1.45, 2.2, 1.15, str(validation.verify_zero_row_count), "验证 0 行", BLUE)
    success_ratio = validation.script_success_count / validation.total_files
    verify_ratio = validation.verify_hit_1_count / validation.total_files
    add_bullet_block(
        slide,
        0.6,
        3.0,
        6.2,
        3.1,
        [
            f"脚本成功率 = {validation.script_success_count}/{validation.total_files} = {success_ratio:.2%}",
            f"验证命中率 = {validation.verify_hit_1_count}/{validation.total_files} = {verify_ratio:.2%}",
            "这里的“成功”定义为 psql 返回码 0，不代表所有验证查询都命中。",
            "唯一的 0-row 成功脚本来自自动命名索引场景：索引创建成功，但 verification query 仍按显式 idx_* 名称查。",
        ],
        title="怎么看这些数字",
        fill_color=RGBColor(30, 41, 59),
        title_color=WHITE,
        body_color=RGBColor(226, 232, 240),
        line_color=RGBColor(51, 65, 85),
    )
    add_code_block(
        slide,
        7.05,
        3.0,
        5.75,
        3.1,
        "psql --version\npg_isready\npsql -d postgres -v ON_ERROR_STOP=1 -f artifacts/generated_sql/...sql\npython3 loop over 1188 files",
        title="本次实跑使用的命令",
    )
    add_footer(slide, "Validation cache: presentation/data/runtime_execution_validation_20260415.json", dark=True)

    # Slide 14
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "13. 失败为什么会集中在 324 个", "失败不是随机噪声，而是 skill 模板里的系统性语义错误")
    error_rows = [[name, str(count)] for name, count in validation.error_counts.items()]
    add_table_slide(slide, 0.55, 1.2, 6.05, 2.05, ["错误类型", "数量"], error_rows)
    add_bullet_block(
        slide,
        0.55,
        3.55,
        6.05,
        2.95,
        [
            "216 个 immutable 错误，精确落在 4 个类型族 × 3 个谓词模板 × 18 个唯一性/排序/NULL 组合。",
            "4 个类型族：temporal、range、multirange、extensible_type。",
            "3 个问题谓词：equality_predicate、range_predicate、expression_predicate_immutable。",
            "说明问题不在随机数据，而在 skill 把 CAST/length 之类表达式直接塞进 index predicate。",
        ],
        title="错误簇 A：谓词不满足 immutable",
    )
    add_bullet_block(
        slide,
        6.85,
        1.2,
        5.95,
        5.3,
        [
            "108 个 operator class 错误，全部来自 geometry 类型族。",
            "精确落在 geometry × 6 个 predicate_mode × 18 个唯一性/排序/NULL 组合。",
            "当前 skill 把 geometry 对应到 point_col，同时 access_method 固定成 btree，PostgreSQL 16.13 没有 point 的默认 btree operator class。",
            "这说明 access_method 与 indexed_type_family 之间缺少兼容性约束。",
        ],
        title="错误簇 B：geometry + btree 不兼容",
    )
    add_footer(slide, "Sources: presentation/data/runtime_execution_validation_20260415.json, sample failures from create_index_*0328.sql and *0750.sql")

    # Slide 15
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "14. 两个典型失败样例", "这两条样例能把失败根因看得很直接")
    add_code_block(
        slide,
        0.55,
        1.2,
        6.0,
        5.55,
        "\n".join(read_text(ROOT / "artifacts" / "generated_sql" / "normal_table-create-index" if False else ROOT / "artifacts" / "generated_sql" / "normal_table-create_index" / "create_index_AC_CREATE_INDEX_000001_0328.sql").splitlines()[76:85]),
        title="失败样例 A：temporal + equality_predicate",
    )
    add_code_block(
        slide,
        6.85,
        1.2,
        5.95,
        5.55,
        "\n".join(read_text(ROOT / "artifacts" / "generated_sql" / "normal_table-create_index" / "create_index_AC_CREATE_INDEX_000001_0750.sql").splitlines()[76:85]),
        title="失败样例 B：geometry(point) + btree",
    )
    add_footer(slide, "Sources: artifacts/generated_sql/normal_table-create_index/create_index_AC_CREATE_INDEX_000001_0328.sql and _0750.sql")

    # Slide 16
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "15. 这套系统做对了什么，又缺了什么", "优点和短板都很明确")
    add_bullet_block(
        slide,
        0.55,
        1.25,
        6.0,
        5.3,
        [
            "优点 1：结构清晰。对象模板、领域规则、流程代码、产物目录边界分明。",
            "优点 2：规则外置。skill YAML 能直接表达因子空间、默认值、模板与绑定。",
            "优点 3：覆盖数学可解释。1188 的来源一眼能算出来。",
            "优点 4：审计结果可追溯。manifest / evaluation / binding.tsv 都是结构化产物。",
            "优点 5：提示词体系分层。通用流程 skill 和领域 DSL 没混在一起。",
        ],
        title="做对了",
    )
    add_bullet_block(
        slide,
        6.8,
        1.25,
        6.0,
        5.3,
        [
            "短板 1：repo 内置 audit 没有真实数据库执行这一层。",
            "短板 2：verification query 没处理 implicit_generated_name 场景。",
            "短板 3：access_method 与类型兼容性没有规则约束。",
            "短板 4：部分 predicate 模板把非 immutable 函数塞进 index predicate。",
            "短板 5：plan TSV 有两条 lifecycle row，但生成 SQL 没兑现 drop_index 分支。",
            "短板 6：rotate_attach 只做到“至少覆盖一次”，非重要因子分布极不均匀。",
        ],
        title="还缺什么",
    )
    add_footer(slide, "Sources: evaluator.py, bindings.tsv, generated SQL, runtime validation")

    # Slide 17
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_title(slide, "16. 如果继续做，这个项目下一步该怎么收敛", "这是我建议直接排进下一轮开发的修正顺序", dark=True)
    add_bullet_block(
        slide,
        0.65,
        1.25,
        12.0,
        5.45,
        [
            "第一优先级：修 predicate_mode 的三个问题模板，替换为 PostgreSQL 明确允许进入 partial index predicate 的 immutable 表达式。",
            "第二优先级：给 geometry / full_text / semi-structured 等类型族补 access method 兼容矩阵，不要把 btree 固定死。",
            "第三优先级：verification query 改成支持显式名和自动名，或在 implicit_generated_name 场景中改查 pg_indexes by tablename + column expression。",
            "第四优先级：把 '执行验证' 变成 workflow 的正式阶段，输出 execution report，而不只是 coverage report。",
            "第五优先级：把 plan TSV 的第二条 lifecycle row 真正落实到脚本生成逻辑里，或者在计划层删除未实现分支。",
            "第六优先级：重新设计 rotate_attach，让非重要因子分布更均匀，避免大量特殊值只出现一次。",
        ],
        title="建议路线",
        fill_color=RGBColor(30, 41, 59),
        title_color=WHITE,
        body_color=RGBColor(226, 232, 240),
        line_color=RGBColor(51, 65, 85),
    )
    add_footer(slide, "Recommendation slide generated from repository findings", dark=True)

    # Slide 18
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "附录 A. 全量因子清单", "这里列出 create_index.skill 中的因子设计")
    add_table_slide(slide, 0.35, 1.1, 12.65, 5.95, ["因子", "标签", "重要性", "取值数", "样例值"], factor_rows)
    add_footer(slide, "Source: skills/index/create_index.skill")

    # Slide 19
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "附录 B. 关键文件索引", "讲项目和提示词时我实际读取了这些文件")
    left_rows = [
        ["README.md", "项目定位"],
        ["PROJECT_STRUCTURE.md", "目录职责"],
        ["src/pg_case_factory/workflow.py", "主流程"],
        ["src/pg_case_factory/request_parser.py", "自然语言约束"],
        ["src/pg_case_factory/abstract_case_generator.py", "展开算法"],
        ["src/pg_case_factory/sql_renderer.py", "SQL 渲染"],
        ["src/pg_case_factory/evaluator.py", "覆盖审计"],
        ["objects/table_02.sql/normal_table/all_column_types.sql", "对象模板"],
    ]
    right_rows = [
        ["skills/index/create_index.skill", "领域规则核心"],
        ["skills/common/output_script_style.skill", "输出脚本规约"],
        ["skills/common/common_factor_policy.skill", "覆盖策略"],
        ["~/.codex/superpowers/skills/using-superpowers/SKILL.md", "skill 入口规约"],
        ["~/.codex/superpowers/skills/brainstorming/SKILL.md", "设计规约"],
        ["~/.codex/superpowers/skills/writing-plans/SKILL.md", "计划规约"],
        ["~/.codex/skills/pg-case-factory-orchestrator/SKILL.md", "repo 专属编排"],
        ["~/.codex/skills/pg-case-factory-plan-auditor/SKILL.md", "repo 专属审计"],
    ]
    add_table_slide(slide, 0.45, 1.3, 6.0, 5.6, ["文件", "用途"], left_rows)
    add_table_slide(slide, 6.8, 1.3, 6.0, 5.6, ["文件", "用途"], right_rows)
    add_footer(slide, "This deck summarizes visible prompt/design files; it does not dump hidden runtime prompt verbatim.")

    # Slide 20
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "附录 C. 精确命令与数据源", "确保这份 PPT 的数字可复核")
    add_code_block(
        slide,
        0.55,
        1.2,
        12.2,
        5.65,
        "\n".join(
            [
                "1. 读取仓库结构与核心代码：sed / rg / find",
                "2. 读取运行产物：artifacts/evaluations/*.json, *.md, manifest.json, bindings.tsv",
                "3. 验证 PostgreSQL 环境：psql --version ; pg_isready",
                "4. 样例执行：psql -d postgres -v ON_ERROR_STOP=1 -f artifacts/generated_sql/..._0001.sql",
                "5. 全量执行：python3 loop over 1188 SQL files and capture psql return code",
                "6. 缓存验证结果：presentation/data/runtime_execution_validation_20260415.json",
                "7. 生成演示文稿：python3 tools/generate_project_ppt.py",
                "",
                "最重要的数据源：",
                f"- {ROOT / 'artifacts' / 'evaluations' / 'TASK_20260414_120051.json'}",
                f"- {ROOT / 'artifacts' / 'intermediates' / 'normal_table-create_index-bindings.tsv'}",
                f"- {ROOT / 'skills' / 'index' / 'create_index.skill'}",
                f"- {ROOT / 'objects' / 'table_02.sql' / 'normal_table' / 'all_column_types.sql'}",
                f"- {VALIDATION_CACHE}",
            ]
        ),
        title="复核路径",
    )
    add_footer(slide, "Generated on 2026-04-15")

    prs.save(OUTPUT_PATH)


if __name__ == "__main__":
    build_deck()
